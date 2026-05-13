"""
make_quad_test.py — generates a test PPTX with a 4-quadrant slide.
Top-left box has "Text Heading" + 3-level dummy bullets for VP-prompt testing.
Run: python make_quad_test.py
Output: quad_test.pptx in the same folder.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

OUT = os.path.join(os.path.dirname(__file__), "quad_test.pptx")

# Slide = 960 x 540 pt  (standard 16:9)
W_SLIDE = Pt(960)
H_SLIDE = Pt(540)

# Quad regions (from LAYOUT_RECIPES.md quad preset), in pt
QUADS = [
    dict(label="TL", left=40,  top=64,  w=432, h=218),
    dict(label="TR", left=488, top=64,  w=432, h=218),
    dict(label="BL", left=40,  top=298, w=432, h=218),
    dict(label="BR", left=488, top=298, w=432, h=218),
]

# Colors (pt-consistent with a financial deck palette)
COLOR_HEADING  = RGBColor(0x1F, 0x38, 0x64)   # dark navy
COLOR_L1       = RGBColor(0x00, 0x00, 0x00)   # black
COLOR_L2       = RGBColor(0x26, 0x4F, 0x82)   # mid-blue
COLOR_L3       = RGBColor(0x70, 0x70, 0x70)   # gray
COLOR_FILL     = RGBColor(0xF2, 0xF2, 0xF2)   # light gray box

# Dummy content per quadrant
# Each item: (indent_level, text)
# indent_level 0 = heading / title line
QUAD_CONTENT = {
    "TL": [
        (0,  "Text Heading"),
        (1,  "Revenue grew 12% YoY to $1.36B"),
        (2,  "EMEA up 18%, APAC up 9%"),
        (2,  "Americas flat at $680M"),
        (1,  "EBITDA margin expanded 120bps to 20.1%"),
        (3,  "Cost discipline drove 60bps"),
        (3,  "Mix shift contributed remaining 60bps"),
        (1,  "Free cash flow of $272M"),
    ],
    "TR": [
        (0,  "Strategic Priorities"),
        (1,  "Accelerate enterprise ARR"),
        (2,  "Top 200 named accounts — land & expand"),
        (1,  "Ship AI copilot to GA"),
        (2,  "Drive attach across installed base"),
        (1,  "Expand gross margin +200bps"),
        (3,  "Multi-cloud cost optimisation"),
    ],
    "BL": [
        (0,  "Risk Factors"),
        (1,  "FX headwinds — 60% revenue offshore"),
        (2,  "Every 1pt USD move = $8M revenue impact"),
        (1,  "Supply chain lead times elevated"),
        (1,  "Regulatory review in EMEA pending"),
        (3,  "Expected resolution Q3 FY25"),
    ],
    "BR": [
        (0,  "Outlook"),
        (1,  "FY25 revenue guidance $1.48–1.52B"),
        (2,  "Assumes stable FX"),
        (1,  "EBITDA margin guide 20.5–21.0%"),
        (1,  "Capex $120–140M"),
        (3,  "Majority datacenter expansion"),
    ],
}

# Bullet char per indent level (level 0 = no bullet / heading)
BULLET_CHARS = {1: "•", 2: "–", 3: "·"}   # •  –  ·

def _set_bullet(para, indent_level):
    """Inject <a:buChar> (or <a:buNone> for heading) into paragraph pPr."""
    pPr = para._p.get_or_add_pPr()
    # Remove any existing bullet elements
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)

    if indent_level == 0:
        buNone = etree.SubElement(pPr, qn("a:buNone"))
    else:
        char = BULLET_CHARS.get(indent_level, "•")
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", char)


def _add_bullet_paragraph(tf, text, indent_level, is_heading=False):
    """Add a paragraph to text frame tf with given indent and formatting."""
    para = tf.add_paragraph()
    para.text = text
    para.level = indent_level
    _set_bullet(para, indent_level)

    run = para.runs[0]
    if is_heading:
        run.font.size  = Pt(13)
        run.font.bold  = True
        run.font.color.rgb = COLOR_HEADING
        para.alignment = PP_ALIGN.LEFT
    else:
        if indent_level == 1:
            run.font.size  = Pt(10)
            run.font.color.rgb = COLOR_L1
        elif indent_level == 2:
            run.font.size  = Pt(10)
            run.font.color.rgb = COLOR_L2
        else:  # level 3+
            run.font.size  = Pt(10)
            run.font.color.rgb = COLOR_L3
        run.font.bold = False

def build():
    prs = Presentation()
    prs.slide_width  = W_SLIDE
    prs.slide_height = H_SLIDE

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    for q in QUADS:
        left  = Pt(q["left"])
        top   = Pt(q["top"])
        width = Pt(q["w"])
        height= Pt(q["h"])

        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.name = f"quad_{q['label']}"

        tf = txBox.text_frame
        tf.word_wrap = True

        # Enable shrink-to-fit via XML (python-pptx doesn't expose this directly)
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("spAutoFit", "0")
            bodyPr.set("normAutofit", "0")
            # shrink-to-fit: set noAutofit=false, add normAutofit element
            # Actually set autofit to normAutofit with fontScale
            for child in list(bodyPr):
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag in ('noAutofit', 'spAutoFit', 'normAutofit'):
                    bodyPr.remove(child)
            norm = etree.SubElement(bodyPr, qn("a:normAutofit"))
            norm.set("fontScale", "100000")
            norm.set("lnSpcReduction", "0")

        # Add box fill via shape XML
        sp = txBox._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(sp, qn("p:spPr"))
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", "F2F2F2")

        content = QUAD_CONTENT[q["label"]]
        for i, (lvl, text) in enumerate(content):
            if i == 0:
                # Use the built-in first paragraph (don't add a new one)
                para = tf.paragraphs[0]
                para.text  = text
                para.level = 0
                _set_bullet(para, 0)
                run = para.runs[0]
                run.font.size      = Pt(13)
                run.font.bold      = True
                run.font.color.rgb = COLOR_HEADING
            else:
                _add_bullet_paragraph(tf, text, lvl, is_heading=False)

    prs.save(OUT)
    print(f"Saved: {OUT}")

if __name__ == "__main__":
    build()
