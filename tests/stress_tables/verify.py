"""
stress_tables/verify.py — validate table advanced ops stress test output.

Opens tests/stress_tables/output/test_tables_v1.pptx via python-pptx (no COM)
and asserts all 6 checks. Exits 0 on full pass, 1 on any failure.

Run: python tests/stress_tables/verify.py
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree

OUT_DECK = Path(__file__).resolve().parent / "output" / "test_tables_v1.pptx"

CHECKS_PASSED = []
CHECKS_FAILED = []

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def check(name: str, condition: bool, detail: str = ""):
    if condition:
        CHECKS_PASSED.append(name)
        print(f"  PASS  {name}")
    else:
        CHECKS_FAILED.append(name)
        print(f"  FAIL  {name}" + (f": {detail}" if detail else ""))


def find_table_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.has_table and shape.name == name:
            return shape
    return None


def get_cell_fill_rgb(cell):
    """Return RGB tuple for a table cell's fill, or None if no explicit fill."""
    try:
        spPr = cell._tc.get_or_add_tcPr()
        solidFill = spPr.find(f"{{{NS}}}solidFill")
        if solidFill is None:
            return None
        srgb = solidFill.find(f"{{{NS}}}srgbClr")
        if srgb is None:
            return None
        val = srgb.get("val", "")
        if len(val) == 6:
            return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
        return None
    except Exception:
        return None


def has_bullet_char(para) -> bool:
    """Return True if paragraph has an explicit buChar element."""
    try:
        pPr = para._p.pPr
        if pPr is None:
            return False
        return pPr.find(f"{{{NS}}}buChar") is not None
    except Exception:
        return False


def get_para_font_color(para) -> tuple:
    """Return (r, g, b) of the first run's explicit font color, or None."""
    try:
        for r in para.runs:
            rgb = r.font.color.rgb
            return (rgb.red, rgb.green, rgb.blue)
    except Exception:
        pass
    return None


def run_checks(prs: Presentation):
    slides = list(prs.slides)
    slide7 = slides[6]   # 0-based index
    slide8 = slides[7]
    slide9 = slides[8]

    # ---- CHECK 1: Slide 7 table has exactly 5 rows and 4 cols ----
    s7_shape = find_table_shape(slide7, "s7_struct")
    if s7_shape is None:
        check("CHECK1_slide7_table_shape_exists", False, "shape 's7_struct' not found on slide 7")
        check("CHECK1b_slide7_table_dimensions", False, "shape not found")
    else:
        check("CHECK1_slide7_table_shape_exists", True)
        tbl = s7_shape.table
        rows_ok = len(tbl.rows) == 5
        cols_ok = len(tbl.columns) == 4
        check("CHECK1b_slide7_table_dimensions",
              rows_ok and cols_ok,
              f"got {len(tbl.rows)} rows x {len(tbl.columns)} cols, expected 5x4")

    # ---- CHECK 2: Slide 7 table cells populated (non-empty content) ----
    populated = False
    if s7_shape is not None:
        tbl = s7_shape.table
        filled = sum(
            1 for row in tbl.rows
            for cell in row.cells
            if cell.text_frame.text.strip()
        )
        populated = filled >= 12  # at least 12 of 20 cells have text
    check("CHECK2_slide7_cells_populated", populated,
          "fewer than 12 of 20 cells have text content")

    # ---- CHECK 3: Slide 8 header row (row 1) has dark-blue fill ----
    s8_shape = find_table_shape(slide8, "s8_format")
    header_fill_ok = False
    if s8_shape is not None:
        tbl = s8_shape.table
        # Check col 1 of row 1 for approx #1F4E79
        cell = tbl.cell(0, 0)
        rgb = get_cell_fill_rgb(cell)
        if rgb:
            r, g, b = rgb
            # Allow ±5 tolerance for color rounding
            header_fill_ok = (abs(r - 0x1F) <= 5 and abs(g - 0x4E) <= 5 and abs(b - 0x79) <= 5)
    check("CHECK3_slide8_header_fill", header_fill_ok,
          "cell(1,1) fill is not close to #1F4E79")

    # ---- CHECK 4: Slide 8 has a merged cell (colspan > 1) ----
    merged_found = False
    if s8_shape is not None:
        tbl = s8_shape.table
        # python-pptx: merged cells have is_spanned=True for non-anchor cells
        # or we look for gridSpan > 1 in the XML
        for row in tbl.rows:
            for cell in row.cells:
                tc = cell._tc
                gridSpan = tc.get("gridSpan", None)
                rowSpan = tc.get("rowSpan", None)
                if (gridSpan is not None and int(gridSpan) > 1) or \
                   (rowSpan is not None and int(rowSpan) > 1):
                    merged_found = True
                    break
            if merged_found:
                break
    check("CHECK4_slide8_merged_cell", merged_found,
          "no merged cell found on slide 8 table")

    # ---- CHECK 5: Slide 9 cell has multiple paragraphs with bullet chars ----
    s9_shape = find_table_shape(slide9, "s9_bullets")
    bullet_cells_found = 0
    if s9_shape is not None:
        tbl = s9_shape.table
        for row in tbl.rows:
            for cell in row.cells:
                paras = cell.text_frame.paragraphs
                bullet_count = sum(1 for p in paras if has_bullet_char(p))
                if bullet_count >= 2:
                    bullet_cells_found += 1
    check("CHECK5_slide9_multi_bullet_cells", bullet_cells_found >= 2,
          f"only {bullet_cells_found} cell(s) have 2+ bullet paragraphs, expected >= 2")

    # ---- CHECK 6: Slide 9 has a cell paragraph with explicit bold + colored text ----
    bold_colored_found = False
    if s9_shape is not None:
        tbl = s9_shape.table
        for row in tbl.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    if not para.runs:
                        continue
                    run = para.runs[0]
                    is_bold = run.font.bold is True
                    try:
                        _ = run.font.color.rgb  # raises if no explicit color
                        has_color = True
                    except Exception:
                        has_color = False
                    if is_bold and has_color:
                        bold_colored_found = True
                        break
                if bold_colored_found:
                    break
            if bold_colored_found:
                break
    check("CHECK6_slide9_bold_colored_paragraph", bold_colored_found,
          "no cell paragraph has both bold=True and explicit font color")


def main() -> int:
    if not OUT_DECK.exists():
        print(f"ERROR: {OUT_DECK} not found — run harness.py first")
        return 1

    prs = Presentation(str(OUT_DECK))
    slides = list(prs.slides)
    if len(slides) < 9:
        print(f"ERROR: deck has only {len(slides)} slides, expected >= 9")
        return 1

    print(f"Verifying {OUT_DECK.name} ({len(slides)} slides)...\n")
    run_checks(prs)

    total = len(CHECKS_PASSED) + len(CHECKS_FAILED)
    print(f"\n{'='*50}")
    print(f"PASSED: {len(CHECKS_PASSED)}/{total}")
    if CHECKS_FAILED:
        print(f"FAILED: {', '.join(CHECKS_FAILED)}")
        return 1
    print("All checks passed.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
