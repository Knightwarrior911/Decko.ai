"""Regenerate test_decks/ from scratch.

Run: python tests/make_test_decks.py
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
DECKS_DIR = REPO_ROOT / "test_decks"


def make_smoke_3slide(path: Path) -> None:
    pres = Presentation()
    layout_title = pres.slide_layouts[0]
    layout_content = pres.slide_layouts[1]

    s1 = pres.slides.add_slide(layout_title)
    s1.shapes.title.text = "Q3 Results"
    s1.placeholders[1].text = "Subtitle text"

    s2 = pres.slides.add_slide(layout_content)
    s2.shapes.title.text = "Bullets"
    s2.placeholders[1].text = "Revenue up 12%\nMargins improved\nHeadcount stable"

    s3 = pres.slides.add_slide(layout_content)
    s3.shapes.title.text = "Next Steps"
    s3.placeholders[1].text = "Plan Q4\nHire 5 engineers"

    pres.save(str(path))


def make_full_visual(path: Path) -> None:
    pres = Presentation()
    layout_blank = pres.slide_layouts[6]

    slide = pres.slides.add_slide(layout_blank)

    # Title textbox
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tb.text_frame
    tf.text = "Q3 Visual Slide"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # Table
    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(2)
    )
    tbl = table_shape.table
    headers = ["Metric", "Q2", "Q3"]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    data = [
        ["Revenue", "100", "112"],
        ["Margin", "26%", "28%"],
    ]
    for r, row in enumerate(data, start=1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v

    # Plain rectangle with fill
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4), Inches(2), Inches(1)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    rect.text_frame.text = "Box"

    pres.save(str(path))


def make_phase2(path: Path) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    pres = Presentation()
    layout_blank = pres.slide_layouts[6]

    # --- Slide 1: multi-paragraph bullet body + speaker notes
    s1 = pres.slides.add_slide(layout_blank)

    title = s1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Bullet Body Slide"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    body = s1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    tf = body.text_frame
    tf.text = "First point about revenue"
    tf.paragraphs[0].runs[0].font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "Second point about margins"
    p2.runs[0].font.size = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Third point about headcount"
    p3.runs[0].font.size = Pt(20)

    s1.notes_slide.notes_text_frame.text = "Highlight Q3 outperformance and Y/Y growth"

    # --- Slide 2: native chart
    s2 = pres.slides.add_slide(layout_blank)
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("FY24", (100, 110, 120, 130))
    s2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.5), Inches(8), Inches(4.5),
        chart_data,
    )

    # --- Slide 3: three ungrouped boxes (group_shapes test will group them)
    s3 = pres.slides.add_slide(layout_blank)
    box1 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(2), Inches(1))
    box1.text_frame.text = "A"
    box2 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.5), Inches(2), Inches(1))
    box2.text_frame.text = "B"
    box3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(2), Inches(1))
    box3.text_frame.text = "C"

    # --- Slide 4: 4x3 table + plain rectangle
    s4 = pres.slides.add_slide(layout_blank)
    rows, cols = 4, 3
    tbl_shape = s4.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(9), Inches(3))
    tbl = tbl_shape.table
    headers = ["Metric", "FY24", "FY25"]
    for c, h in enumerate(headers):
        tbl.cell(0, c).text = h
    body_data = [
        ["Revenue", "100", "112"],
        ["Margin", "26%", "28%"],
        ["Headcount", "1,200", "1,250"],
    ]
    for r, row in enumerate(body_data, start=1):
        for c, v in enumerate(row):
            tbl.cell(r, c).text = v

    rect = s4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(2), Inches(0.8)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    rect.text_frame.text = "Plain"

    pres.save(str(path))


def make_text_v3(path):
    """Fixture for granular text v3:
    - Slide 1: heading + paragraph with mixed-format runs
    - Slide 2: bullets w/ mixed sizes + one hyperlink + one strikethrough
    - Slide 3: paragraphs with sub/superscript
    - Slide 4: shape with non-default text frame (vertical-mid, margins, line-spacing 1.5)
    """
    import win32com.client
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True
    pres = app.Presentations.Add(WithWindow=True)
    try:
        # Slide 1: heading + mixed-format paragraph
        s1 = pres.Slides.Add(1, 1)  # ppLayoutText
        s1.Shapes.Title.TextFrame.TextRange.Text = "Earnings Q3"
        body = s1.Shapes(2).TextFrame.TextRange
        body.Text = "Revenue grew 23% in Q3"
        # Bold "grew 23%"
        bold_range = body.Characters(9, 8)  # "grew 23%"
        bold_range.Font.Bold = -1

        # Slide 2: bullets, mixed sizes, hyperlink, strikethrough
        s2 = pres.Slides.Add(2, 2)  # ppLayoutBullet
        s2.Shapes.Title.TextFrame.TextRange.Text = "Highlights"
        bullets = s2.Shapes(2).TextFrame.TextRange
        bullets.Text = ("First point about revenue\r"
                        "Second point with a link\r"
                        "Third item, deprecated")
        # First bullet: size 24 on the word "revenue"
        bullets.Paragraphs(1).Characters(20, 7).Font.Size = 24
        # Second bullet: hyperlink on "link"
        link_range = bullets.Paragraphs(2).Characters(22, 4)
        link_range.ActionSettings(1).Hyperlink.Address = "https://decko.ai/docs"
        # Third bullet: strikethrough on "deprecated" — applied via XML patch after SaveAs
        # (PowerPoint COM typelib v2.12 does not expose Font.Strikethrough)

        # Slide 3: sub/superscript
        s3 = pres.Slides.Add(3, 1)
        s3.Shapes.Title.TextFrame.TextRange.Text = "Chemistry"
        body3 = s3.Shapes(2).TextFrame.TextRange
        body3.Text = "H2O\rE=mc2"
        body3.Paragraphs(1).Characters(2, 1).Font.BaselineOffset = -0.25  # sub on "2"
        body3.Paragraphs(2).Characters(5, 1).Font.BaselineOffset = 0.30   # super on "2"

        # Slide 4: text-frame anchored middle, margins, line-spacing 1.5
        s4 = pres.Slides.Add(4, 1)
        s4.Shapes.Title.TextFrame.TextRange.Text = "Layout"
        sh4 = s4.Shapes(2)
        tf4 = sh4.TextFrame
        tf4.TextRange.Text = "Anchored middle\rWith line spacing"
        tf4.VerticalAnchor = 3  # msoAnchorMiddle
        tf4.MarginLeft   = 18.0
        tf4.MarginRight  = 18.0
        tf4.MarginTop    = 9.0
        tf4.MarginBottom = 9.0
        for i in range(1, 3):
            p = tf4.TextRange.Paragraphs(i).ParagraphFormat
            p.LineRuleWithin = -1   # msoTrue
            p.SpaceWithin    = 1.5

        pres.SaveAs(str(path), 24)  # ppSaveAsOpenXMLPresentation
    finally:
        pres.Close()
        app.Quit()

    # Post-process: patch strikethrough on "deprecated" via python-pptx XML
    # COM typelib v2.12 does not expose Font.Strikethrough, so we apply it here.
    from pptx import Presentation as _Pres
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree

    _p = _Pres(str(path))
    # Slide 2 (index 1), shape index 1 (body placeholder), paragraph index 2 (3rd bullet)
    _tf = _p.slides[1].placeholders[1].text_frame
    # The third paragraph text is "Third item, deprecated"
    # chars 13..22 = "deprecated" (1-based in COM = 0-based index 12..21)
    # We need to split the run to isolate "deprecated" and add strike
    _para = _tf.paragraphs[2]
    # Get the text of the paragraph to locate "deprecated"
    _full_text = _para.text  # "Third item, deprecated"
    _strike_start = _full_text.find("deprecated")
    assert _strike_start >= 0, f"'deprecated' not found in: {_full_text!r}"
    _strike_end = _strike_start + len("deprecated")
    # Rebuild runs: before | "deprecated" (struck) | after
    _before = _full_text[:_strike_start]   # "Third item, "
    _struck = _full_text[_strike_start:_strike_end]  # "deprecated"
    _after  = _full_text[_strike_end:]     # ""
    _pXml = _para._p
    # Remove all existing <a:r> children
    for _r in _pXml.findall(_qn("a:r")):
        _pXml.remove(_r)
    # Add run for text before "deprecated"
    if _before:
        _r1 = _etree.SubElement(_pXml, _qn("a:r"))
        _t1 = _etree.SubElement(_r1, _qn("a:t"))
        _t1.text = _before
    # Add struck run for "deprecated"
    _r2 = _etree.SubElement(_pXml, _qn("a:r"))
    _rpr2 = _etree.SubElement(_r2, _qn("a:rPr"), attrib={"lang": "en-US", "strike": "sngStrike", "dirty": "0"})
    _t2 = _etree.SubElement(_r2, _qn("a:t"))
    _t2.text = _struck
    # Add run for text after
    if _after:
        _r3 = _etree.SubElement(_pXml, _qn("a:r"))
        _t3 = _etree.SubElement(_r3, _qn("a:t"))
        _t3.text = _after
    _p.save(str(path))


def main() -> int:
    DECKS_DIR.mkdir(parents=True, exist_ok=True)
    make_smoke_3slide(DECKS_DIR / "smoke_3slide.pptx")
    make_full_visual(DECKS_DIR / "full_visual.pptx")
    make_phase2(DECKS_DIR / "phase2.pptx")
    make_text_v3(DECKS_DIR / "text_v3.pptx")
    print("text_v3.pptx OK")
    print(f"[done] wrote 4 decks to {DECKS_DIR}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
