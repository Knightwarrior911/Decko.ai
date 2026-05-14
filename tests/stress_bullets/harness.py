"""
stress_bullets/harness.py — bullet formatting stress test via Decko COM.

Copies JAZZ-Pitch-Book.pptx (never modifies original), then replaces slides 4-6
with 3 stress-test slides covering:
  Slide 4: text box — no-bullet header + 5 bullets (3 indent levels, 3 styles,
            mixed bold+color within a paragraph via add_run)
  Slide 5: rrect shape — same multi-level structure + one all-bold paragraph
  Slide 6: 3x2 table — bullets in cells, cell para with level>0 via set_cell_indent_level

Run: python tests/stress_bullets/harness.py
"""

import json
import shutil
import sys
from pathlib import Path

import win32com.client
from pptx import Presentation
from lxml import etree

REPO_ROOT = Path(__file__).resolve().parent.parent.parent
CARRIER   = REPO_ROOT / "PPT_AI_Editor.pptm"
JAZZ_SRC  = Path(r"C:\Users\vinit\Downloads\JAZZ-Pitch-Book.pptx")
OUT_DIR   = Path(__file__).resolve().parent / "output"
OUT_DECK  = OUT_DIR / "test_bullets_v2.pptx"

OUT_DIR.mkdir(exist_ok=True)


def run_batch(app, actions: list) -> str:
    instr = json.dumps({"actions": actions, "verify_after": False}, ensure_ascii=True)
    return str(app.Run("PPT_AI_Editor.pptm!modExecuteInstructions.ExecuteFromString", instr))


def slide4_actions() -> list:
    """Text box: no-bullet header + 5 paragraphs across 3 bullet levels, mixed run formatting."""
    S = 4
    REF = "s4_textbox"

    # --- PHASE 1: structural (clear + add shape + add ALL paragraphs) ---
    # Para 0 (via add_text_box): "Performance Highlights"  <- no-bullet heading
    # Para 1: "Revenue grew 18% YoY"         <- level 0 disc (will split runs later)
    # Para 2: "North America up 22% driven"   <- level 1 dash (will split runs later)
    # Para 3: "Europe flat; macro headwinds"  <- level 1 dash
    # Para 4: "UK down 3%, Germany up 8%"     <- level 2 square
    # Para 5: "EBITDA margin expanded 210bps" <- level 0 disc (will split runs later)

    acts = [
        {"type": "clear_slide", "slide": S},
        {"type": "add_text_box", "slide": S, "ref_name": REF,
         "text": "Performance Highlights",
         "pos": {"left": 50, "top": 40, "width": 620, "height": 440},
         "font_size": 11},
        # Add paragraphs in order BEFORE any formatting
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 0,
         "value": "Revenue grew 18% YoY"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 1,
         "value": "North America up 22% driven by pharma"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 2,
         "value": "Europe flat; macro headwinds persist"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 3,
         "value": "UK down 3%, Germany up 8%"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 4,
         "value": "EBITDA margin expanded 210bps"},
    ]

    # --- PHASE 2: paragraph-level formatting (bullet style + indent level) ---
    acts += [
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 0, "value": "none"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 1, "value": "disc"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 2, "value": "dash"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 3, "value": "dash"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 4, "value": "square"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 5, "value": "disc"},
        {"type": "set_indent_level", "slide": S, "shape_id": REF, "paragraph_index": 2, "value": 1},
        {"type": "set_indent_level", "slide": S, "shape_id": REF, "paragraph_index": 3, "value": 1},
        {"type": "set_indent_level", "slide": S, "shape_id": REF, "paragraph_index": 4, "value": 2},
    ]

    # --- PHASE 3: run-level mixed formatting via add_run ---
    # Para 1: shorten run 0 → "Revenue grew ", add bold+colored run, add plain run
    acts += [
        {"type": "set_run_text", "slide": S, "shape_id": REF,
         "paragraph_index": 1, "run_index": 0, "value": "Revenue grew "},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 1, "value": "18% YoY", "bold": True, "color": "#C00000"},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 1, "value": " in FY2024"},
    ]
    # Para 2: shorten run 0 → "North America up ", add bold run, add plain, add bold+colored
    acts += [
        {"type": "set_run_text", "slide": S, "shape_id": REF,
         "paragraph_index": 2, "run_index": 0, "value": "North America up "},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 2, "value": "22%", "bold": True},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 2, "value": " driven by "},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 2, "value": "pharma", "bold": True, "color": "#1F4E79"},
    ]
    # Para 5: shorten run 0, add bold run
    acts += [
        {"type": "set_run_text", "slide": S, "shape_id": REF,
         "paragraph_index": 5, "run_index": 0, "value": "EBITDA margin expanded "},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 5, "value": "210bps", "bold": True},
    ]

    return acts


def slide5_actions() -> list:
    """Rounded rect shape with same multi-level bullets + one all-bold paragraph."""
    S = 5
    REF = "s5_rrect"

    acts = [
        {"type": "clear_slide", "slide": S},
        {"type": "add_shape", "slide": S, "kind": "rrect", "ref_name": REF,
         "pos": {"left": 50, "top": 40, "width": 620, "height": 440},
         "fill": "#F2F2F2", "stroke": "#CCCCCC", "stroke_weight_pt": 0.75},
        # Para 0 via set_text: no-bullet heading
        {"type": "set_text", "slide": S, "shape_id": REF, "value": "Strategic Priorities"},
        # Add paragraphs
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 0,
         "value": "Expand North America market share"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 1,
         "value": "Target pharma and biotech verticals"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 2,
         "value": "Key accounts: Pfizer, Roche, Novartis"},
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 3,
         "value": "Divest underperforming EU assets"},
        # All-bold paragraph: entire paragraph bold (tests all-bold = single bold run)
        {"type": "add_paragraph", "slide": S, "shape_id": REF, "after_paragraph_index": 4,
         "value": "CRITICAL: Board approval required Q1 2025"},
    ]

    # Paragraph formatting
    acts += [
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 0, "value": "none"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 1, "value": "disc"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 2, "value": "dash"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 3, "value": "square"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 4, "value": "disc"},
        {"type": "set_bullet_style", "slide": S, "shape_id": REF, "paragraph_index": 5, "value": "disc"},
        {"type": "set_indent_level", "slide": S, "shape_id": REF, "paragraph_index": 2, "value": 1},
        {"type": "set_indent_level", "slide": S, "shape_id": REF, "paragraph_index": 3, "value": 2},
    ]

    # All-bold paragraph (para 5): set run 0 bold
    acts += [
        {"type": "set_run_bold", "slide": S, "shape_id": REF, "paragraph_index": 5, "run_index": 0, "value": True},
        {"type": "set_run_font_color", "slide": S, "shape_id": REF, "paragraph_index": 5, "run_index": 0, "value": "#7030A0"},
    ]

    # Mixed runs in para 1
    acts += [
        {"type": "set_run_text", "slide": S, "shape_id": REF,
         "paragraph_index": 1, "run_index": 0, "value": "Expand "},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 1, "value": "North America", "bold": True, "color": "#1F4E79"},
        {"type": "add_run", "slide": S, "shape_id": REF,
         "paragraph_index": 1, "value": " market share"},
    ]

    return acts


def slide6_actions() -> list:
    """3x2 table: bullets in cells, cell para at level>0 via set_cell_indent_level."""
    S = 6
    REF = "s6_table"

    acts = [
        {"type": "clear_slide", "slide": S},
        {"type": "add_table", "slide": S, "ref_name": REF, "rows": 3, "cols": 2,
         "pos": {"left": 50, "top": 40, "width": 620, "height": 350}},
    ]

    # Cell (1,1): no-bullet header + 2 bullet paragraphs
    acts += [
        {"type": "set_cell_text", "slide": S, "shape_id": REF, "row": 1, "col": 1, "value": "Pipeline"},
        {"type": "add_cell_paragraph", "slide": S, "shape_id": REF, "row": 1, "col": 1,
         "after_paragraph_index": 0, "value": "Phase 3 trials: 4 assets"},
        {"type": "add_cell_paragraph", "slide": S, "shape_id": REF, "row": 1, "col": 1,
         "after_paragraph_index": 1, "value": "Approvals expected 2025-2026"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 1, "col": 1, "paragraph_index": 0, "value": "none"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 1, "col": 1, "paragraph_index": 1, "value": "disc"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 1, "col": 1, "paragraph_index": 2, "value": "disc"},
    ]

    # Cell (1,2): header + 2 bullets including level>0 via set_cell_indent_level
    acts += [
        {"type": "set_cell_text", "slide": S, "shape_id": REF, "row": 1, "col": 2, "value": "Financials"},
        {"type": "add_cell_paragraph", "slide": S, "shape_id": REF, "row": 1, "col": 2,
         "after_paragraph_index": 0, "value": "Revenue: $480M (+18%)"},
        {"type": "add_cell_paragraph", "slide": S, "shape_id": REF, "row": 1, "col": 2,
         "after_paragraph_index": 1, "value": "EBITDA: $96M (20% margin)"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 1, "col": 2, "paragraph_index": 0, "value": "none"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 1, "col": 2, "paragraph_index": 1, "value": "disc"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 1, "col": 2, "paragraph_index": 2, "value": "disc"},
        # NOTE: set_cell_indent_level intentionally omitted here.
        # PowerPoint COM cannot set table cell paragraph levels via ParagraphFormat.IndentLevel.
        # Handled in Phase 2 (python-pptx) after Decko saves.
    ]

    # Cell (2,1): colored text
    acts += [
        {"type": "set_cell_text", "slide": S, "shape_id": REF, "row": 2, "col": 1,
         "value": "Key Risks"},
        {"type": "add_cell_paragraph", "slide": S, "shape_id": REF, "row": 2, "col": 1,
         "after_paragraph_index": 0, "value": "Regulatory delay risk: HIGH"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 2, "col": 1, "paragraph_index": 0, "value": "none"},
        {"type": "set_cell_bullet_style", "slide": S, "shape_id": REF,
         "row": 2, "col": 1, "paragraph_index": 1, "value": "disc"},
        {"type": "set_cell_paragraph_font_color", "slide": S, "shape_id": REF,
         "row": 2, "col": 1, "paragraph_index": 1, "value": "#C00000"},
    ]

    return acts


NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def phase2_set_cell_para_levels(deck_path: Path):
    """Phase 2: set cell paragraph indent levels via python-pptx (VBA cannot do this).
    Sets slide 6, table shape 'REF s6_table', cell (1,2) paragraph index 2 to level 1.
    """
    prs = Presentation(str(deck_path))
    slides = list(prs.slides)
    slide6 = slides[5]

    for shape in slide6.shapes:
        if shape.has_table and shape.name == "s6_table":
            cell = shape.table.cell(0, 1)   # (row=1,col=2) → 0-indexed = (0,1)
            paras = cell.text_frame.paragraphs
            if len(paras) > 2:
                # Set paragraph index 2 (3rd paragraph) to level 1
                p_elem = paras[2]._p
                pPr = p_elem.get_or_add_pPr()
                pPr.set("lvl", "1")
                print("  Phase 2: set cell(1,2) para[2] lvl=1 via python-pptx")
            break

    prs.save(str(deck_path))


def main() -> int:
    if not JAZZ_SRC.exists():
        print(f"ERROR: {JAZZ_SRC} not found")
        return 1
    if not CARRIER.exists():
        print(f"ERROR: {CARRIER} not found")
        return 1

    print(f"Copying {JAZZ_SRC.name} -> {OUT_DECK}")
    shutil.copy2(JAZZ_SRC, OUT_DECK)

    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True

    while app.Presentations.Count > 0:
        try:
            app.Presentations(1).Close()
        except Exception:
            break

    print("Opening carrier...")
    carrier = app.Presentations.Open(str(CARRIER), WithWindow=False)

    print("Opening test deck...")
    deck = app.Presentations.Open(str(OUT_DECK), WithWindow=True)
    deck.Windows(1).Activate()

    for slide_num, build_fn in [(4, slide4_actions), (5, slide5_actions), (6, slide6_actions)]:
        print(f"Building slide {slide_num}...")
        acts = build_fn()
        result = run_batch(app, acts)
        applied = [x for x in result.split(",") if "applied" in x]
        skipped = [x for x in result.split(",") if "skipped" in x or "error" in x.lower()]
        print(f"  Slide {slide_num}: {result.strip()[:120]}")
        if skipped:
            print(f"  WARNING: {skipped}")

    print("Saving deck...")
    deck.SaveAs(str(OUT_DECK))
    deck.Close()
    carrier.Saved = True
    carrier.Close()
    app.Quit()

    print("Phase 2: setting cell paragraph levels via python-pptx...")
    phase2_set_cell_para_levels(OUT_DECK)

    print(f"Done. Output: {OUT_DECK}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
