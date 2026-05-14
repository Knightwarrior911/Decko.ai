"""
stress_bullets/verify.py — validate bullet formatting stress test output.

Opens tests/stress_bullets/output/test_bullets.pptx via python-pptx (no COM)
and asserts all 6 checks. Exits 0 on full pass, 1 on any failure.

Run: python tests/stress_bullets/verify.py
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

OUT_DECK = Path(__file__).resolve().parent / "output" / "test_bullets_v2.pptx"

CHECKS_PASSED = []
CHECKS_FAILED = []


def check(name: str, condition: bool, detail: str = ""):
    if condition:
        CHECKS_PASSED.append(name)
        print(f"  PASS  {name}")
    else:
        CHECKS_FAILED.append(name)
        print(f"  FAIL  {name}" + (f": {detail}" if detail else ""))


def bullet_type(para) -> str:
    """Return the bullet type string for a paragraph."""
    try:
        pPr = para._p.pPr
        if pPr is None:
            return "none"
        buNone = pPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
        )
        if buNone is not None:
            return "none"
        buChar = pPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
        )
        if buChar is not None:
            return f"char:{buChar.get('char', '?')}"
        buAutoNum = pPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum"
        )
        if buAutoNum is not None:
            return "autonumber"
        return "inherited"
    except Exception:
        return "unknown"


def has_explicit_bullet(para) -> bool:
    bt = bullet_type(para)
    return bt not in ("none", "inherited", "unknown") and not bt.startswith("none")


def run_checks(prs: Presentation):
    slides = list(prs.slides)
    slide4 = slides[3]   # 0-based
    slide5 = slides[4]
    slide6 = slides[5]

    # ---- CHECK 1: 3 distinct indent levels in one text frame ----
    found_3_levels = False
    for slide in [slide4, slide5]:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            levels = {p.level for p in shape.text_frame.paragraphs if p.text.strip()}
            if len(levels) >= 3:
                found_3_levels = True
                break
        if found_3_levels:
            break
    check("CHECK1_three_indent_levels", found_3_levels,
          "no text frame with paragraphs at levels 0, 1, 2")

    # ---- CHECK 2: 2+ different bullet styles in one text frame ----
    found_mixed_styles = False
    for slide in [slide4, slide5]:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            styles = set()
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    styles.add(bullet_type(p))
            if len(styles) >= 2:
                found_mixed_styles = True
                break
        if found_mixed_styles:
            break
    check("CHECK2_mixed_bullet_styles", found_mixed_styles,
          "no text frame has 2+ distinct bullet styles")

    # ---- CHECK 3: paragraph[0] no-bullet, paragraph[1] has bullet ----
    found_header_then_bullet = False
    for slide in [slide4, slide5, slide6]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                paras = [p for p in shape.text_frame.paragraphs]
            elif shape.has_table:
                # Check table cells too
                tbl = shape.table
                for row in tbl.rows:
                    for cell in row.cells:
                        cparas = [p for p in cell.text_frame.paragraphs]
                        if (len(cparas) >= 2
                                and not has_explicit_bullet(cparas[0])
                                and has_explicit_bullet(cparas[1])):
                            found_header_then_bullet = True
                continue
            else:
                continue
            if (len(paras) >= 2
                    and not has_explicit_bullet(paras[0])
                    and has_explicit_bullet(paras[1])):
                found_header_then_bullet = True
                break
        if found_header_then_bullet:
            break
    check("CHECK3_header_then_bullet", found_header_then_bullet,
          "no shape/cell found with para[0]=no-bullet and para[1]=bullet")

    # ---- CHECK 4: mixed bold runs within one paragraph ----
    found_mixed_bold = False
    for slide in [slide4, slide5]:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if len(runs) < 2:
                    continue
                bolds = [r.font.bold for r in runs]
                true_bold  = any(b is True  for b in bolds)
                false_bold = any(b is False or b is None for b in bolds)
                if true_bold and false_bold:
                    found_mixed_bold = True
                    break
            if found_mixed_bold:
                break
        if found_mixed_bold:
            break
    check("CHECK4_mixed_bold_runs", found_mixed_bold,
          "no paragraph has both bold=True run and non-bold run")

    # ---- CHECK 5: mixed colored vs uncolored runs in one paragraph ----
    found_mixed_color = False
    from pptx.dml.color import RGBColor
    for slide in [slide4, slide5]:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if len(runs) < 2:
                    continue
                colored = []
                for r in runs:
                    try:
                        c = r.font.color.rgb
                        colored.append(True)
                    except Exception:
                        colored.append(False)
                if any(colored) and not all(colored):
                    found_mixed_color = True
                    break
            if found_mixed_color:
                break
        if found_mixed_color:
            break
    check("CHECK5_mixed_color_runs", found_mixed_color,
          "no paragraph has both a colored run and an uncolored run")

    # ---- CHECK 6: table cell paragraph with level > 0 ----
    found_cell_indent = False
    for shape in slide6.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        for row in tbl.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    if para.level > 0:
                        found_cell_indent = True
                        break
                if found_cell_indent:
                    break
            if found_cell_indent:
                break
        if found_cell_indent:
            break
    check("CHECK6_table_cell_indent_level", found_cell_indent,
          "no table cell paragraph has level > 0")


def main() -> int:
    if not OUT_DECK.exists():
        print(f"ERROR: {OUT_DECK} not found — run harness.py first")
        return 1

    prs = Presentation(str(OUT_DECK))
    slides = list(prs.slides)
    if len(slides) < 6:
        print(f"ERROR: deck has only {len(slides)} slides, expected >= 6")
        return 1

    print(f"Verifying {OUT_DECK.name} ({len(slides)} slides)...\n")
    run_checks(prs)

    print(f"\n{'='*50}")
    print(f"PASSED: {len(CHECKS_PASSED)}/6")
    if CHECKS_FAILED:
        print(f"FAILED: {', '.join(CHECKS_FAILED)}")
        return 1
    print("All 6 checks passed.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
