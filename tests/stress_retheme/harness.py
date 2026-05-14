"""
stress_retheme/harness.py — stress-test the recolor_deck action on JAZZ-Pitch-Book.pptx.

Steps:
  1. Pre-scan JAZZ for top-4 explicit colors (excluding black/white).
  2. Save sidecar colors.json.
  3. Scenario A: remap top-4 -> PALETTE_A (4 mappings in one call).
  4. Scenario B: on Scenario A output, remap #003087->#1F4E79, #D03027->#C00000.
  5. Scenario C: fresh JAZZ copy, remap top-1 -> #FFFFFF.

Exit 0 if all 3 scenarios have 0 skipped actions.
Run: python tests/stress_retheme/harness.py
"""

import sys, json, re, shutil, time
from pathlib import Path
from collections import Counter

ROOT     = Path(__file__).resolve().parents[2]
JAZZ     = Path(r"C:\Users\vinit\Downloads\JAZZ-Pitch-Book.pptx")
OUT_DIR  = Path(__file__).resolve().parent / "output"
SIDECAR  = OUT_DIR / "colors.json"
CARRIER  = ROOT / "PPT_AI_Editor.pptm"

PALETTE_A = ["#003087", "#D03027", "#F5F5F5", "#15283C"]

SCENARIO_A = OUT_DIR / "jazz_retheme_A.pptx"
SCENARIO_B = OUT_DIR / "jazz_retheme_B.pptx"
SCENARIO_C = OUT_DIR / "jazz_retheme_C.pptx"

SKIP_TRIVIAL = {"#000000", "#FFFFFF", "#FEFEFE", "#FDFDFD"}


# ── color scanning (python-pptx, no COM) ─────────────────────────────────────

def scan_colors(pptx_path: Path) -> Counter:
    from pptx import Presentation
    counter: Counter = Counter()
    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        _scan_shapes(slide.shapes, counter)
    return counter


def _scan_shapes(shapes, counter: Counter) -> None:
    for sh in shapes:
        try:
            rgb = sh.fill.fore_color.rgb
            counter[f"#{str(rgb).upper()}"] += 1
        except Exception:
            pass
        try:
            rgb = sh.line.color.rgb
            counter[f"#{str(rgb).upper()}"] += 1
        except Exception:
            pass
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rgb = run.font.color.rgb
                        counter[f"#{str(rgb).upper()}"] += 1
                    except Exception:
                        pass
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    try:
                        rgb = cell.fill.fore_color.rgb
                        counter[f"#{str(rgb).upper()}"] += 1
                    except Exception:
                        pass
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                rgb = run.font.color.rgb
                                counter[f"#{str(rgb).upper()}"] += 1
                            except Exception:
                                pass
        if sh.shape_type == 6:
            try:
                _scan_shapes(sh.shapes, counter)
            except Exception:
                pass


# ── COM helpers ───────────────────────────────────────────────────────────────

def open_and_activate(app, path: Path):
    """Open presentation, activate its window, return prs object."""
    prs = app.Presentations.Open(str(path.resolve()), WithWindow=True)
    time.sleep(1.5)
    try:
        prs.Windows(1).Activate()
        time.sleep(0.5)
    except Exception:
        pass
    return prs


def run_actions(app, prs, actions: list) -> str:
    payload = json.dumps({"actions": actions, "verify_after": False}, ensure_ascii=True)
    result = app.Run(
        "PPT_AI_Editor.pptm!modExecuteInstructions.ExecuteFromString",
        payload
    )
    return str(result)


def check_skipped(result_str: str, scenario: str) -> bool:
    m = re.search(r'(\d+) applied.*?(\d+) skipped', result_str)
    if not m:
        print(f"  WARN  [{scenario}] cannot parse result: {result_str[:300]}")
        return False
    applied, skipped = int(m.group(1)), int(m.group(2))
    if skipped > 0:
        print(f"  FAIL  [{scenario}] {skipped} skipped. Result: {result_str[:400]}")
        return False
    print(f"  OK    [{scenario}] {applied} applied, 0 skipped")
    return True


# ── main ──────────────────────────────────────────────────────────────────────

def main() -> int:
    if not JAZZ.exists():
        print(f"ERROR: JAZZ deck not found: {JAZZ}")
        return 1
    if not CARRIER.exists():
        print(f"ERROR: Carrier not found: {CARRIER}")
        return 1

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    # ── Step 1: scan for top-4 colors ────────────────────────────────────────
    print(f"Scanning {JAZZ.name} for explicit colors...")
    counter = scan_colors(JAZZ)
    top4 = [c for c, _ in counter.most_common(30) if c not in SKIP_TRIVIAL][:4]
    if len(top4) < 4:
        top4 = [c for c, _ in counter.most_common(4)]
    print(f"  Top-4 (non-trivial): {top4}")

    SIDECAR.write_text(json.dumps(
        {"found_colors": top4, "counts": {c: counter[c] for c in top4}},
        indent=2
    ))
    print(f"  Sidecar: {SIDECAR}")

    # ── Step 2: build action payloads ─────────────────────────────────────────
    mappings_A = [{"from": top4[i], "to": PALETTE_A[i]} for i in range(4)]
    mappings_B = [
        {"from": "#003087", "to": "#1F4E79"},
        {"from": "#D03027", "to": "#C00000"},
    ]
    mappings_C = [{"from": top4[0], "to": "#FFFFFF"}]

    actions_A = [{"type": "recolor_deck", "mappings": mappings_A}]
    actions_B = [{"type": "recolor_deck", "mappings": mappings_B}]
    actions_C = [{"type": "recolor_deck", "mappings": mappings_C}]

    # ── Step 3: COM runs ──────────────────────────────────────────────────────
    import win32com.client
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True

    # close any stale presentations from prior runs
    while app.Presentations.Count > 0:
        try:
            app.Presentations(1).Close()
        except Exception:
            break

    carrier = app.Presentations.Open(str(CARRIER.resolve()), WithWindow=False)
    time.sleep(2)

    all_ok = True

    try:
        # ── Scenario A ───────────────────────────────────────────────────────
        print(f"\nScenario A: remap {top4} -> {PALETTE_A}")
        work_a = OUT_DIR / "jazz_a_work.pptx"
        shutil.copy2(str(JAZZ), str(work_a))
        prs_a = open_and_activate(app, work_a)
        res_a = run_actions(app, prs_a, actions_A)
        prs_a.SaveAs(str(SCENARIO_A.resolve()), 24)
        prs_a.Close()
        all_ok = check_skipped(res_a, "A") and all_ok
        work_a.unlink(missing_ok=True)

        # ── Scenario B ───────────────────────────────────────────────────────
        print("\nScenario B: remap #003087->#1F4E79, #D03027->#C00000 (on A output)")
        work_b = OUT_DIR / "jazz_b_work.pptx"
        shutil.copy2(str(SCENARIO_A), str(work_b))
        prs_b = open_and_activate(app, work_b)
        res_b = run_actions(app, prs_b, actions_B)
        prs_b.SaveAs(str(SCENARIO_B.resolve()), 24)
        prs_b.Close()
        all_ok = check_skipped(res_b, "B") and all_ok
        work_b.unlink(missing_ok=True)

        # ── Scenario C ───────────────────────────────────────────────────────
        print(f"\nScenario C: remap {top4[0]} -> #FFFFFF (fresh JAZZ copy)")
        work_c = OUT_DIR / "jazz_c_work.pptx"
        shutil.copy2(str(JAZZ), str(work_c))
        prs_c = open_and_activate(app, work_c)
        res_c = run_actions(app, prs_c, actions_C)
        prs_c.SaveAs(str(SCENARIO_C.resolve()), 24)
        prs_c.Close()
        all_ok = check_skipped(res_c, "C") and all_ok
        work_c.unlink(missing_ok=True)

    finally:
        # close all non-carrier presentations
        try:
            for i in range(app.Presentations.Count, 0, -1):
                try:
                    prs_i = app.Presentations(i)
                    if prs_i.FullName.lower() != str(CARRIER.resolve()).lower():
                        prs_i.Close()
                except Exception:
                    pass
        except Exception:
            pass
        try:
            carrier.Close()
        except Exception:
            pass
        try:
            app.Quit()
        except Exception:
            pass

    if all_ok:
        print("\nAll 3 scenarios: 0 skipped actions.")
        return 0
    print("\nFAIL: some scenarios had skipped actions.")
    return 1


if __name__ == "__main__":
    sys.exit(main())
