"""
stress_retheme/verify.py — validate recolor_deck stress test output.

Opens output/*.pptx via python-pptx (no COM) and asserts 6 checks.
Exits 0 on full pass, 1 on any failure.

Run: python tests/stress_retheme/verify.py
"""

import sys, json
from pathlib import Path
from collections import Counter
from pptx import Presentation

OUT_DIR   = Path(__file__).resolve().parent / "output"
JAZZ_ORIG = Path(r"C:\Users\vinit\Downloads\JAZZ-Pitch-Book.pptx")
SIDECAR   = OUT_DIR / "colors.json"
SCENARIO_A = OUT_DIR / "jazz_retheme_A.pptx"
SCENARIO_B = OUT_DIR / "jazz_retheme_B.pptx"
SCENARIO_C = OUT_DIR / "jazz_retheme_C.pptx"

CHECKS_PASSED: list = []
CHECKS_FAILED: list = []


def check(name: str, condition: bool, detail: str = "") -> None:
    if condition:
        CHECKS_PASSED.append(name)
        print(f"  PASS  {name}")
    else:
        CHECKS_FAILED.append(name)
        print(f"  FAIL  {name}" + (f": {detail}" if detail else ""))


def collect_colors(pptx_path: Path) -> set:
    """Return set of '#RRGGBB' explicit colors found in the deck."""
    colors: set = set()
    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        _scan(slide.shapes, colors)
    return colors


def _scan(shapes, colors: set) -> None:
    for sh in shapes:
        try:
            colors.add(f"#{str(sh.fill.fore_color.rgb).upper()}")
        except Exception:
            pass
        try:
            colors.add(f"#{str(sh.line.color.rgb).upper()}")
        except Exception:
            pass
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        colors.add(f"#{str(run.font.color.rgb).upper()}")
                    except Exception:
                        pass
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    try:
                        colors.add(f"#{str(cell.fill.fore_color.rgb).upper()}")
                    except Exception:
                        pass
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                colors.add(f"#{str(run.font.color.rgb).upper()}")
                            except Exception:
                                pass
        if sh.shape_type == 6:
            try:
                _scan(sh.shapes, colors)
            except Exception:
                pass


def run_checks() -> None:
    # ── CHECK 1: all 3 output files exist ────────────────────────────────────
    all_exist = SCENARIO_A.exists() and SCENARIO_B.exists() and SCENARIO_C.exists()
    missing = [p.name for p in [SCENARIO_A, SCENARIO_B, SCENARIO_C] if not p.exists()]
    check("CHECK1_all_outputs_exist", all_exist,
          f"missing: {missing}" if missing else "")

    if not all_exist:
        for name in ["CHECK2", "CHECK3", "CHECK4", "CHECK5", "CHECK6"]:
            check(f"{name}_skipped_missing_outputs", False, "outputs missing")
        return

    # ── CHECK 2: Scenario A has same slide count as JAZZ original ─────────────
    jazz_slides = len(list(Presentation(str(JAZZ_ORIG)).slides)) if JAZZ_ORIG.exists() else 51
    a_slides = len(list(Presentation(str(SCENARIO_A)).slides))
    check("CHECK2_scenario_A_slide_count",
          a_slides == jazz_slides,
          f"got {a_slides}, expected {jazz_slides}")

    # ── CHECK 3: sidecar has 4 found_colors ───────────────────────────────────
    if SIDECAR.exists():
        sidecar = json.loads(SIDECAR.read_text())
        found = sidecar.get("found_colors", [])
        check("CHECK3_sidecar_has_4_colors",
              len(found) == 4,
              f"got {len(found)} colors: {found}")
    else:
        check("CHECK3_sidecar_has_4_colors", False, "colors.json not found")

    # ── CHECK 4: Scenario A contains #003087 (top-1 original remapped) ────────
    print("  Scanning Scenario A colors...")
    colors_a = collect_colors(SCENARIO_A)
    check("CHECK4_scenario_A_has_003087",
          "#003087" in colors_a,
          f"#003087 not found; sample colors: {sorted(colors_a)[:10]}")

    # ── CHECK 5: Scenario B contains #1F4E79 (secondary remap of #003087) ─────
    print("  Scanning Scenario B colors...")
    colors_b = collect_colors(SCENARIO_B)
    check("CHECK5_scenario_B_has_1F4E79",
          "#1F4E79" in colors_b,
          f"#1F4E79 not found; sample colors: {sorted(colors_b)[:10]}")

    # ── CHECK 6: Scenario A also contains #D03027 (multi-mapping verified) ────
    check("CHECK6_scenario_A_has_D03027",
          "#D03027" in colors_a,
          f"#D03027 not found; sample colors: {sorted(colors_a)[:10]}")


def main() -> int:
    print("Verifying stress_retheme outputs...\n")
    run_checks()

    total = len(CHECKS_PASSED) + len(CHECKS_FAILED)
    print(f"\n{'='*50}")
    print(f"PASSED: {len(CHECKS_PASSED)}/{total}")
    if CHECKS_FAILED:
        print(f"FAILED: {', '.join(CHECKS_FAILED)}")
        return 1
    print("All checks passed.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
